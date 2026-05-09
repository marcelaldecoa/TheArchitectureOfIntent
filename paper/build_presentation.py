"""
Build a teaching deck for the paper.

Output: paper/architecture-of-intent.pptx
Design: 16:9, dark navy accents, two figures embedded from paper/figures/.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


HERE = Path(__file__).resolve().parent
FIG_TREE = HERE / "figures" / "archetype-decision-tree.png"
FIG_ORTHO = HERE / "figures" / "four-dimensions-orthogonality.png"
OUT = HERE / "architecture-of-intent.pptx"


# Palette
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
ACCENT = RGBColor(0x1E, 0x3A, 0x8A)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PALE = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def setup_widescreen(prs):
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=NAVY, font="Calibri", bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, body = item
        else:
            head, body = item, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "■  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        bullet.font.name = font
        run = p.add_run()
        run.text = head
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = font
        if body:
            sub = p.add_run()
            sub.text = "  " + body
            sub.font.size = Pt(size - 2)
            sub.font.bold = False
            sub.font.color.rgb = SLATE
            sub.font.name = font
    return tb


def add_rect(slide, left, top, width, height, fill=PALE, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_top_bar(slide, label):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.45), fill=NAVY)
    add_text(slide, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.4),
             label, size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page, total):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.35),
             "The Architecture of Intent  —  Aldecoa, 2026",
             size=10, color=SLATE)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35),
             f"{page} / {total}", size=10, color=SLATE, align=PP_ALIGN.RIGHT)


# =========================================================================

prs = Presentation()
setup_widescreen(prs)

slides_meta = []  # (build_fn, section_label)


def slide(label):
    def deco(fn):
        slides_meta.append((fn, label))
        return fn
    return deco


# -------------------------------------------------------------------------
@slide("Title")
def s_title(s):
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=NAVY)
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.15), Inches(2.3), fill=AMBER)
    add_text(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(0.6),
             "THE ARCHITECTURE OF INTENT", size=44, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(3.2), Inches(11.5), Inches(0.6),
             "A Framework for Designing Delegated Systems",
             size=24, color=PALE, italic=True)
    add_text(s, Inches(1.0), Inches(4.4), Inches(11.5), Inches(0.4),
             "Marcel Aldecoa  ·  2026",
             size=18, color=PALE)
    add_text(s, Inches(1.0), Inches(4.9), Inches(11.5), Inches(0.4),
             "Position-and-framework paper  ·  ~15,000 words / 34 pages",
             size=14, color=PALE)
    add_text(s, Inches(1.0), Inches(6.6), Inches(11.5), Inches(0.4),
             "Companion teaching deck",
             size=12, italic=True, color=PALE)


# -------------------------------------------------------------------------
@slide("Roadmap")
def s_roadmap(s):
    add_top_bar(s, "ROADMAP")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "What we'll cover today",
             size=32, bold=True, color=NAVY)
    items = [
        ("1. The problem",
         "Why intent is now a design surface, and the judgment gap that opened in 2024–2026"),
        ("2. The framework",
         "Five archetypes · four calibration dimensions · seven failure categories · SDD as protocol"),
        ("3. Application",
         "Coding agents, computer-use agents, and composition with Microsoft DevSquad Copilot"),
        ("4. Honest contribution accounting",
         "What is novel, what is borrowed, and what the framework does NOT claim"),
        ("5. Limitations & take-homes",
         "Where the framework breaks, what to read next, what to try Monday morning"),
    ]
    add_bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.3),
                items, size=20)


# -------------------------------------------------------------------------
@slide("Setup")
def s_problem(s):
    add_top_bar(s, "1. THE PROBLEM")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Delegation is a structural problem, not a model problem",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
             "Humans must express intent precisely enough that a non-human executor can act on it without supervisory rescue.",
             size=16, italic=True, color=SLATE)
    items = [
        ("Software", "spec → compiler → binary"),
        ("Organizations", "policy → manager → individual contributor"),
        ("Automated pipelines", "config → scheduler → worker"),
        ("AI agents", "prompt + tools → model → action  —  the most acute current instance"),
    ]
    add_bullets(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(3.5),
                items, size=20)
    add_rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0), fill=PALE,
             line=AMBER)
    add_text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.8),
             "Claim: intent is a primary design artifact distinct from implementation. AI agents are\nthe class where this claim becomes load-bearing — because the executor is probabilistic.",
             size=16, italic=True, color=NAVY)


# -------------------------------------------------------------------------
@slide("Setup")
def s_judgment_gap(s):
    add_top_bar(s, "1. THE PROBLEM")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "The judgment gap (2024–2026)",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.6),
             "An execution-first bias in current agent harnesses",
             size=18, italic=True, color=SLATE)
    add_rect(s, Inches(0.6), Inches(2.4), Inches(5.9), Inches(4.4),
             fill=PALE, line=ACCENT)
    add_text(s, Inches(0.85), Inches(2.55), Inches(5.4), Inches(0.5),
             "Senior engineer", size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.85), Inches(3.1), Inches(5.4), Inches(3.6),
                [("Surfaces ambiguity", "before acting"),
                 ("Asks \"is this what you mean?\"", ""),
                 ("Refuses ill-posed work", ""),
                 ("Treats clarification", "as load-bearing")],
                size=15)
    add_rect(s, Inches(6.85), Inches(2.4), Inches(5.9), Inches(4.4),
             fill=PALE, line=AMBER)
    add_text(s, Inches(7.1), Inches(2.55), Inches(5.4), Inches(0.5),
             "LLM agent (even with AskUserQuestion)",
             size=18, bold=True, color=AMBER)
    add_bullets(s, Inches(7.1), Inches(3.1), Inches(5.4), Inches(3.6),
                [("Executes first", "even on under-specified work"),
                 ("Pauses rarely", "even when tools are available"),
                 ("Reaches for action", "before reaching for clarification"),
                 ("Tool exists; judgment", "to use it is not yet calibrated")],
                size=15)


# -------------------------------------------------------------------------
@slide("Setup")
def s_claims(s):
    add_top_bar(s, "1. THE PROBLEM")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Honest contribution accounting",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Reviewers reward narrow novelty claims. Overclaiming gets punished.",
             size=16, italic=True, color=SLATE)
    add_rect(s, Inches(0.6), Inches(2.4), Inches(5.9), Inches(4.4),
             fill=PALE, line=ACCENT)
    add_text(s, Inches(0.85), Inches(2.55), Inches(5.4), Inches(0.5),
             "What IS novel (3 things)",
             size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.85), Inches(3.1), Inches(5.4), Inches(3.6),
                [("Cat 7 (Perceptual Failure)", "for perceiving-then-acting agents"),
                 ("Orthogonality of agency × autonomy", "extending Shavit & Agarwal 2023"),
                 ("Fix-locus framing", "of the failure taxonomy")],
                size=15)
    add_rect(s, Inches(6.85), Inches(2.4), Inches(5.9), Inches(4.4),
             fill=PALE, line=SLATE)
    add_text(s, Inches(7.1), Inches(2.55), Inches(5.4), Inches(0.5),
             "What is NOT claimed as novel",
             size=18, bold=True, color=SLATE)
    add_bullets(s, Inches(7.1), Inches(3.1), Inches(5.4), Inches(3.6),
                [("SDD as discipline", "lineage: spec-kit, DevSquad"),
                 ("Archetypes as concept", "lineage: Anthropic Building Effective Agents"),
                 ("The four dimensions individually", "lineage: SAE J3016, Shavit & Agarwal"),
                 ("Cat 1–6", "synthesis from common practice")],
                size=14)


# -------------------------------------------------------------------------
@slide("Framework")
def s_framework_overview(s):
    add_top_bar(s, "2. THE FRAMEWORK")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Four load-bearing elements",
             size=32, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "The synthesis is the larger contribution. The four bind into a framework you can apply in design, in spec review, and in post-incident diagnosis.",
             size=14, italic=True, color=SLATE)
    cards = [
        ("Archetypes", "Five canonical delegation shapes\nAdvisor · Executor · Guardian\nSynthesizer · Orchestrator", ACCENT),
        ("Calibration dimensions", "Four orthogonal dials\nAgency · Autonomy\nResponsibility · Reversibility", ACCENT),
        ("Failure taxonomy", "Seven categories by fix locus\nCat 1 (Spec) → Cat 6 (Model)\nCat 7 (Perceptual) — NOVEL", AMBER),
        ("Spec-Driven Development", "The executable protocol\nthat binds the above into\na spec the agent runs", ACCENT),
    ]
    x = Inches(0.6)
    w = Inches(3.0)
    gap = Inches(0.13)
    for i, (head, body, color) in enumerate(cards):
        left = Inches(0.6 + i * (3.0 + 0.13))
        add_rect(s, left, Inches(2.4), w, Inches(4.3), fill=PALE, line=color)
        add_rect(s, left, Inches(2.4), w, Inches(0.55), fill=color)
        add_text(s, left + Inches(0.15), Inches(2.45), w - Inches(0.3), Inches(0.45),
                 head, size=15, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.2), Inches(3.15), w - Inches(0.4), Inches(3.4),
                 body, size=14, color=NAVY)


# -------------------------------------------------------------------------
@slide("Framework")
def s_archetypes(s):
    add_top_bar(s, "2.1 ARCHETYPES")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Five canonical delegation shapes",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Pre-commit a system to one shape before design begins. Composition is permitted.",
             size=15, italic=True, color=SLATE)
    rows = [
        ("Advisor",      "Surfaces options; the human acts",      "Code-review suggester, search agent"),
        ("Executor",     "Acts within a defined scope",            "CI/CD pipeline, deploy agent"),
        ("Guardian",     "Vetoes; protects a boundary",             "Compliance gate, safety filter"),
        ("Synthesizer",  "Combines inputs into a new whole",        "Research summarizer, report generator"),
        ("Orchestrator", "Coordinates other agents/services",       "Multi-agent supervisor, planner"),
    ]
    headers = ["Archetype", "Primary intent", "Examples"]
    col_x = [Inches(0.6), Inches(3.4), Inches(8.0)]
    col_w = [Inches(2.7), Inches(4.5), Inches(5.0)]
    add_rect(s, Inches(0.6), Inches(2.4), Inches(12.0), Inches(0.5), fill=NAVY)
    for i, h in enumerate(headers):
        add_text(s, col_x[i] + Inches(0.1), Inches(2.4), col_w[i],
                 Inches(0.5), h, size=15, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(rows):
        y = Inches(2.95 + r * 0.78)
        bg = PALE if r % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), y, Inches(12.0), Inches(0.78), fill=bg)
        for i, val in enumerate(row):
            bold = (i == 0)
            color = ACCENT if i == 0 else NAVY
            add_text(s, col_x[i] + Inches(0.1), y, col_w[i], Inches(0.78),
                     val, size=14, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)


# -------------------------------------------------------------------------
@slide("Framework")
def s_decision_tree(s):
    add_top_bar(s, "2.1 ARCHETYPES")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "The archetype selection tree",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Apply questions in order. Stop at the first match. Risk override applies.",
             size=15, italic=True, color=SLATE)
    if FIG_TREE.exists():
        s.shapes.add_picture(str(FIG_TREE),
                             Inches(3.4), Inches(2.0),
                             height=Inches(5.0))
    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "Figure 1 (paper). Composition is permitted: a deployment may host multiple archetypes.",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)


# -------------------------------------------------------------------------
@slide("Framework")
def s_dimensions(s):
    add_top_bar(s, "2.2 CALIBRATION DIMENSIONS")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Four orthogonal dials, set deliberately in the spec",
             size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Most teams collapse these into \"how autonomous is the agent?\" — that intuition hides the calibration work.",
             size=14, italic=True, color=SLATE)
    cells = [
        ("Agency", "discretion", "How much judgment when instructions don't fully cover the situation"),
        ("Autonomy", "operation", "How much of the work runs without human intervention at each step"),
        ("Responsibility", "accountability", "How accountability is distributed across the humans around the agent"),
        ("Reversibility", "consequence", "How easy or hard it is to undo what the agent did"),
    ]
    for i, (name, axis, body) in enumerate(cells):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.25)
        y = Inches(2.4 + row * 2.3)
        add_rect(s, x, y, Inches(6.0), Inches(2.0), fill=PALE, line=ACCENT)
        add_rect(s, x, y, Inches(6.0), Inches(0.5), fill=ACCENT)
        add_text(s, x + Inches(0.2), y, Inches(5.6), Inches(0.5),
                 name + "  —  the " + axis + " dimension",
                 size=16, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), y + Inches(0.65), Inches(5.5), Inches(1.3),
                 body, size=15, color=NAVY)


# -------------------------------------------------------------------------
@slide("Framework")
def s_orthogonality(s):
    add_top_bar(s, "2.2 CALIBRATION DIMENSIONS")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Agency × Autonomy: orthogonal, not collinear",
             size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Treating agency and autonomy as a single \"automation level\" (e.g., SAE J3016) collapses this design space onto a diagonal. Real systems sit in all four quadrants.",
             size=13, italic=True, color=SLATE)
    if FIG_ORTHO.exists():
        s.shapes.add_picture(str(FIG_ORTHO),
                             Inches(2.6), Inches(2.4),
                             height=Inches(4.5))
    add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
             "Figure 2 (paper). Responsibility and reversibility are similarly orthogonal.",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)


# -------------------------------------------------------------------------
@slide("Framework")
def s_taxonomy(s):
    add_top_bar(s, "2.3 FAILURE TAXONOMY")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Seven categories, organized by fix locus",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "The diagnostic test: \"If a competent agent had executed this spec as written, would the outcome have been correct?\"",
             size=14, italic=True, color=SLATE)
    rows = [
        ("Cat 1", "Spec",        "Spec was incomplete, ambiguous, or wrong",   "Update the spec"),
        ("Cat 2", "Capability",  "Agent lacked or misused a tool",              "Add or fix the tool / manifest"),
        ("Cat 3", "Scope creep", "Agent did adjacent work it wasn't authorized","Tighten NOT-authorized clauses"),
        ("Cat 4", "Oversight",   "Error escaped the review checkpoint",          "Redesign the oversight model"),
        ("Cat 5", "Compounding", "Early error cascaded through later steps",     "Checkpoint at the handoff"),
        ("Cat 6", "Model-level", "Model failed despite a correct spec",          "Narrow scope, switch model, or accept residual risk"),
        ("Cat 7", "Perceptual",  "Perception diverged from environment state",   "Structural controls + verification step  (NOVEL)"),
    ]
    col_x = [Inches(0.6), Inches(1.4), Inches(3.2), Inches(7.2)]
    col_w = [Inches(0.8), Inches(1.7), Inches(4.0), Inches(5.4)]
    headers = ["#", "Name", "Failure shape", "Fix locus"]
    add_rect(s, Inches(0.6), Inches(2.3), Inches(12.0), Inches(0.45), fill=NAVY)
    for i, h in enumerate(headers):
        add_text(s, col_x[i] + Inches(0.08), Inches(2.3), col_w[i],
                 Inches(0.45), h, size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(rows):
        y = Inches(2.75 + r * 0.62)
        is_cat7 = (r == 6)
        bg = RGBColor(0xFE, 0xF3, 0xC7) if is_cat7 else (PALE if r % 2 == 0 else WHITE)
        add_rect(s, Inches(0.6), y, Inches(12.0), Inches(0.62), fill=bg)
        for i, val in enumerate(row):
            bold = (i <= 1)
            color = AMBER if (is_cat7 and i <= 1) else (ACCENT if i <= 1 else NAVY)
            add_text(s, col_x[i] + Inches(0.08), y, col_w[i], Inches(0.62),
                     val, size=12, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)


# -------------------------------------------------------------------------
@slide("Framework")
def s_cat7(s):
    add_top_bar(s, "2.3 FAILURE TAXONOMY  ·  THE NOVEL CATEGORY")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Cat 7 — Perceptual Failure",
             size=30, bold=True, color=AMBER)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "The system's perception of the environment diverged from the actual state, and the system acted on the wrong perception.",
             size=15, italic=True, color=SLATE)
    add_text(s, Inches(0.6), Inches(2.3), Inches(6.0), Inches(0.4),
             "Four sub-categories", size=18, bold=True, color=NAVY)
    sub = [
        ("Misidentification", "Confirmation gate"),
        ("Missed element",    "Screenshot-then-verify"),
        ("Hallucinated element", "Element-allowlist + DOM grounding"),
        ("State miscount",    "Re-verify at action time"),
    ]
    for i, (name, fix) in enumerate(sub):
        y = Inches(2.85 + i * 0.85)
        add_rect(s, Inches(0.6), y, Inches(6.0), Inches(0.75),
                 fill=PALE, line=AMBER)
        add_text(s, Inches(0.85), y, Inches(3.5), Inches(0.75),
                 name, size=15, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.3), y, Inches(2.2), Inches(0.75),
                 fix, size=13, color=AMBER,
                 anchor=MSO_ANCHOR.MIDDLE, italic=True)
    add_text(s, Inches(7.0), Inches(2.3), Inches(6.0), Inches(0.4),
             "Why a separate category", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(2.85), Inches(6.0), Inches(4.0),
                [("Prior taxonomies don't partition this",
                  "MAST, hallucination survey, OWASP LLM Top 10"),
                 ("It applies only to perceiving-then-acting systems",
                  "computer-use · browser-use · robotic"),
                 ("The fixes don't live in the prompt",
                  "they live in the spec, the manifest, and the verification step"),
                 ("Smallest contribution that gives the taxonomy a slot",
                  "for a class of failure already observed in deployments")],
                size=13)


# -------------------------------------------------------------------------
@slide("Framework")
def s_sdd(s):
    add_top_bar(s, "2.4 SPEC-DRIVEN DEVELOPMENT")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "SDD: the executable protocol layer",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "A spec the agent can run and a human can validate. Lineage: spec-kit (GitHub), DevSquad (Microsoft).",
             size=14, italic=True, color=SLATE)
    add_text(s, Inches(0.6), Inches(2.3), Inches(6.0), Inches(0.4),
             "13 canonical sections", size=18, bold=True, color=NAVY)
    sections = [
        "1. Problem statement", "2. Objective",
        "3. Authorized scope", "4. NOT-authorized scope",
        "5. Tool manifest", "6. Invariants",
        "7. Non-functional constraints", "8. Acceptance criteria",
        "9. Agent execution instructions", "10. Oversight model",
        "11. Validation checklist", "12. Spec evolution log",
        "13. Spec gap log",
    ]
    for i, sec in enumerate(sections):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 3.2)
        y = Inches(2.85 + row * 0.5)
        add_text(s, x, y, Inches(3.0), Inches(0.4),
                 "•  " + sec, size=12, color=NAVY)
    add_text(s, Inches(7.2), Inches(2.3), Inches(6.0), Inches(0.4),
             "How the framework maps in", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(7.2), Inches(2.85), Inches(5.8), Inches(4.0),
                [("Agency lives in", "§3 (scope) and §4 (NOT-authorized)"),
                 ("Autonomy lives in", "§10 (oversight model)"),
                 ("Responsibility lives in", "§1 (problem) and §11 (validation)"),
                 ("Reversibility lives in", "§5 (manifest) and §6 (invariants)"),
                 ("Each failure category", "maps to a specific update site")],
                size=13)


# -------------------------------------------------------------------------
@slide("Application")
def s_coding_agents(s):
    add_top_bar(s, "3. APPLICATION  ·  CODING AGENTS")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Coding agents — where structural fixes compound",
             size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Most coding-agent failures look novel; nearly all map to a Cat 1–6 fix-locus.",
             size=14, italic=True, color=SLATE)
    items = [
        ("The deleted-tests failure",
         "Cat 1 / Cat 3 hybrid — fix in the spec (§4 NOT-authorized) and in CI (test-skip set monotonic), never only in the prompt"),
        ("Three structural controls",
         "(1) PR-only branch protection  (2) tool manifest excludes unrestricted shell  (3) test-suite invariant in CI"),
        ("The discipline",
         "structural fixes live in spec / manifest / CI / platform — never only in the prompt"),
        ("What Cat 7 does NOT do here",
         "coding agents are text-only; they have no perception–action interface that can diverge from environment state"),
    ]
    add_bullets(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(4.6),
                items, size=16)


# -------------------------------------------------------------------------
@slide("Application")
def s_computer_use(s):
    add_top_bar(s, "3. APPLICATION  ·  COMPUTER-USE AGENTS")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Computer-use — where Cat 7 becomes load-bearing",
             size=26, bold=True, color=AMBER)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Anthropic Computer Use · OpenAI Operator · Gemini computer use — vision-then-action systems.",
             size=14, italic=True, color=SLATE)
    add_text(s, Inches(0.6), Inches(2.3), Inches(6.0), Inches(0.4),
             "New failure surfaces", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(2.85), Inches(6.0), Inches(4.0),
                [("Lookalike-domain navigation", "homoglyph / typo / subdomain confusion"),
                 ("Visual instruction injection", "rendered text treated as authoritative"),
                 ("Modal popup interception", "adversarial dialog as legitimate prompt")],
                size=14)
    add_text(s, Inches(7.0), Inches(2.3), Inches(6.0), Inches(0.4),
             "Four structural controls", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(2.85), Inches(6.0), Inches(4.0),
                [("Sandboxed environment", "no host credentials, files, extensions"),
                 ("Authentication scope minimization", "redirects don't carry session cookies"),
                 ("Domain allowlist", "non-allowlisted redirect halts and surfaces"),
                 ("High-consequence confirmation gates", "irreversible-action class always gates")],
                size=14)


# -------------------------------------------------------------------------
@slide("Application")
def s_devsquad(s):
    add_top_bar(s, "3. APPLICATION  ·  DEVSQUAD COMPOSITION")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Composition with Microsoft DevSquad Copilot (2026)",
             size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "DevSquad provides the cadence; the framework provides the design vocabulary. They compose cleanly when the artifact mapping is explicit.",
             size=14, italic=True, color=SLATE)
    rows = [
        ("Phase 1–2", "Envisioning + Spec the slice", "Archetype committed; provisional dimensions"),
        ("Phase 3",       "Plan only what the slice needs",   "Invariants in §6; ADRs constrain §8"),
        ("Phase 4",       "Decompose that slice",           "Least-Capability tool subset per task"),
        ("Phase 5",       "Implement with TDD discipline",  "Spec acceptance suite gates each commit"),
        ("Phase 6",       "Learn in the open",               "Categorize Cat 1–7; trace to fix locus"),
        ("Phase 7",       "Review in independent context",  "Validation in fresh sub-agent context"),
        ("Phase 8",       "Refine continuously",            "Four signal metrics drive next sprint"),
    ]
    headers = ["Phase", "DevSquad activity", "Framework artifact / discipline"]
    col_x = [Inches(0.6), Inches(2.0), Inches(6.6)]
    col_w = [Inches(1.4), Inches(4.5), Inches(6.0)]
    add_rect(s, Inches(0.6), Inches(2.4), Inches(12.0), Inches(0.45), fill=NAVY)
    for i, h in enumerate(headers):
        add_text(s, col_x[i] + Inches(0.08), Inches(2.4), col_w[i],
                 Inches(0.45), h, size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(rows):
        y = Inches(2.85 + r * 0.6)
        bg = PALE if r % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), y, Inches(12.0), Inches(0.6), fill=bg)
        for i, val in enumerate(row):
            bold = (i == 0)
            color = ACCENT if i == 0 else NAVY
            add_text(s, col_x[i] + Inches(0.08), y, col_w[i], Inches(0.6),
                     val, size=12, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)


# -------------------------------------------------------------------------
@slide("Closing")
def s_limits(s):
    add_top_bar(s, "5. LIMITATIONS")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "Where this framework breaks",
             size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Position-and-framework paper. No quantitative validation at scale yet — acknowledged explicitly.",
             size=14, italic=True, color=SLATE)
    items = [
        ("Position paper, not empirical",
         "Applied across three worked examples; quantitative validation across many independent deployments is future work."),
        ("Coarse partitions",
         "Five archetypes, four dimensions, seven categories — opinionated rather than derived."),
        ("Multi-tenant fleet governance is missing",
         "The framework currently does not address the governance of large fleets of agents at scale."),
        ("Cardinality is judgment, not theorem",
         "Reasonable readers may propose six categories or eight; the framework's working choice is seven."),
        ("Cat 7 sub-categories may evolve",
         "Computer-use is young; the four sub-categories will likely be revised as the deployment surface matures."),
    ]
    add_bullets(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(4.6),
                items, size=15)


# -------------------------------------------------------------------------
@slide("Closing")
def s_takehomes(s):
    add_top_bar(s, "5. TAKE-HOMES")
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             "What to try Monday morning",
             size=28, bold=True, color=NAVY)
    items = [
        ("Pick one agent system you're operating",
         "Run the archetype decision tree against it. Does the result match what you have today?"),
        ("Write the spec for its most consequential action",
         "Use the canonical template. Focus on §3 (scope) and §4 (NOT-authorized). The first draft will feel incomplete — that is the point."),
        ("Categorize the first three failures",
         "Apply the diagnostic test, walk Cat 1–7, trace to the fix locus. How many were spec gaps versus model limits?"),
        ("Surface ambiguity, don't resolve",
         "When the spec is unclear, the agent's job is to surface, not to guess. Wire AskUserQuestion into the harness AND calibrate the judgment to use it."),
        ("Compound the learning",
         "Each diagnosed failure improves a spec or a skill. The improvement is durable and propagates to every future task."),
    ]
    add_bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.3),
                items, size=15)


# -------------------------------------------------------------------------
@slide("Closing")
def s_qna(s):
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=NAVY)
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.15), Inches(2.3), fill=AMBER)
    add_text(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(0.6),
             "Discussion", size=44, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(3.3), Inches(11.5), Inches(0.6),
             "Three prompts to start the conversation",
             size=20, italic=True, color=PALE)
    qs = [
        "1.  Where in your team is intent currently implicit, and where would making it explicit pay off most?",
        "2.  What's the last agent failure you saw — and which Cat 1–7 was it, really?",
        "3.  Are your structural controls (spec, manifest, CI, platform) doing the work that today still lives in prompts?",
    ]
    for i, q in enumerate(qs):
        add_text(s, Inches(1.0), Inches(4.4 + i * 0.7), Inches(11.5), Inches(0.6),
                 q, size=18, color=PALE)
    add_text(s, Inches(1.0), Inches(6.7), Inches(11.5), Inches(0.4),
             "Paper: marcelaldecoa.github.io/TheArchitectureOfIntent  ·  Companion book: same site",
             size=12, italic=True, color=PALE)


# =========================================================================
# Build all slides
total = len(slides_meta)
for i, (build_fn, label) in enumerate(slides_meta, start=1):
    s = blank(prs)
    build_fn(s)
    if i > 1 and i < total:
        add_footer(s, i, total)

prs.save(OUT)
print(f"Wrote {OUT} ({total} slides)")
